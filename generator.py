"""
generator.py — Professional PowerPoint & HTML Generator
Consistent fonts, colours, and layout inspired by real-world
academic / corporate presentation standards.
"""

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt, Emu
import os
import base64
import html as html_lib

# ---------------------------------------------------------------------------
# Design System — colours, fonts, spacing
# ---------------------------------------------------------------------------

FONT_TITLE   = "Segoe UI"     # Modern, authoritative branding font
FONT_BODY    = "Segoe UI"     # Clean, readable body font
FONT_FOOTER  = "Segoe UI"

# Slide dimensions  (widescreen 16:9 = 10" × 5.625")
SLIDE_W = Inches(10)
SLIDE_H = Inches(5.625)

THEMES = {
    "blue": {
        "primary"  : RGBColor(31,  73,  125),   # Deep navy blue
        "secondary": RGBColor(68, 114, 196),     # Lighter blue
        "accent"   : RGBColor(255, 192,   0),    # Gold accent
        "bg"       : RGBColor(245, 248, 255),    # Near-white tinted
        "text"     : RGBColor(30,  30,  30),
        "footer"   : RGBColor(100, 100, 120),
    },
    "red": {
        "primary"  : RGBColor(180,  30,  30),
        "secondary": RGBColor(220,  60,  60),
        "accent"   : RGBColor(255, 200,   0),
        "bg"       : RGBColor(255, 248, 248),
        "text"     : RGBColor(30,  30,  30),
        "footer"   : RGBColor(120, 80,  80),
    },
    "green": {
        "primary"  : RGBColor(14,  100,  55),
        "secondary": RGBColor(39,  174,  96),
        "accent"   : RGBColor(243, 156,  18),
        "bg"       : RGBColor(245, 255, 250),
        "text"     : RGBColor(20,  40,  20),
        "footer"   : RGBColor(80,  120,  80),
    },
    "purple": {
        "primary"  : RGBColor(102,  51, 153),
        "secondary": RGBColor(155,  89, 182),
        "accent"   : RGBColor(230, 126,  34),
        "bg"       : RGBColor(248, 245, 255),
        "text"     : RGBColor(30,  20,  50),
        "footer"   : RGBColor(100,  80, 130),
    },
    "dark": {
        "primary"  : RGBColor(30,  144, 255),   # Dodger blue on dark
        "secondary": RGBColor(70,  130, 180),
        "accent"   : RGBColor(255, 215,   0),
        "bg"       : RGBColor(28,   28,  36),   # Very dark navy
        "text"     : RGBColor(230, 230, 240),
        "footer"   : RGBColor(140, 140, 160),
    },
    "premium_gold": {
        "primary"  : RGBColor(184, 134,  11),   # Dark Goldenrod
        "secondary": RGBColor(218, 165,  32),   # Goldenrod
        "accent"   : RGBColor(255, 255, 255),
        "bg"       : RGBColor(18,   18,  18),   # Obsidian
        "text"     : RGBColor(245, 245, 245),
        "footer"   : RGBColor(100, 100, 100),
    },
}

WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(0,   0,   0)


def _rgb(theme: dict, key: str) -> RGBColor:
    return theme.get(key, RGBColor(50, 50, 50))


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _add_textbox(slide, left, top, width, height,
                 text: str, font_name: str, font_size: Pt,
                 bold: bool = False, color: RGBColor = BLACK,
                 align=PP_ALIGN.LEFT, wrap: bool = True):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.text      = text
    p.alignment = align
    run = p.runs[0] if p.runs else p.add_run()
    run.text            = text
    run.font.name       = font_name
    run.font.size       = font_size
    run.font.bold       = bold
    run.font.color.rgb  = color
    return txb


def _add_rect(slide, left, top, width, height, fill: RGBColor, line=False):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line:
        shape.line.color.rgb = fill
    else:
        shape.line.fill.background()
    return shape


def _auto_font_size(text: str, base: int = 44,
                    thresholds=((60, 32), (100, 26), (160, 20), (300, 16))) -> Pt:
    n = len(text)
    for limit, size in thresholds:
        if n <= limit:
            return Pt(size)
    return Pt(base - 28)   # smallest fallback


# ---------------------------------------------------------------------------
# Cover Slide
# ---------------------------------------------------------------------------

def _build_cover_slide(prs, cover_title: str, student_name: str,
                       student_id: str, theme: dict, cover_image_path=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])   # blank layout

    W, H = prs.slide_width, prs.slide_height

    if cover_image_path and os.path.exists(cover_image_path):
        # PDF cover as full-bleed background
        slide.shapes.add_picture(cover_image_path, 0, 0, W, H)
        # Dark overlay for readability
        overlay = _add_rect(slide, 0, 0, W, H, RGBColor(10, 20, 40))
        # Manually set alpha-like effect via brightness
        overlay.fill.fore_color.rgb = RGBColor(10, 20, 40)
        text_color = WHITE
        accent_bar = _rgb(theme, "accent")
    else:
        # Solid branded background
        _add_rect(slide, 0, 0, W, H, _rgb(theme, "bg"))
        # Left accent strip
        _add_rect(slide, 0, 0, Inches(0.35), H, _rgb(theme, "primary"))
        text_color = _rgb(theme, "text")
        accent_bar = _rgb(theme, "primary")

    # Horizontal accent bar (top 8% of slide height)
    bar_h = Inches(0.55)
    _add_rect(slide, 0, 0, W, bar_h, _rgb(theme, "primary"))

    # Gold / accent bottom bar
    bottom_bar_h = Inches(0.45)
    _add_rect(slide, 0, H - bottom_bar_h, W, bottom_bar_h, _rgb(theme, "accent"))

    # --- Title text box (centred vertically) ---
    title_top  = Inches(1.2)
    title_h    = Inches(3.0)
    title_left = Inches(0.7)
    title_w    = W - Inches(1.4)

    txb = slide.shapes.add_textbox(title_left, title_top, title_w, title_h)
    tf  = txb.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text           = cover_title
    run.font.name      = FONT_TITLE
    run.font.size      = _auto_font_size(cover_title)
    run.font.bold      = True
    run.font.color.rgb = WHITE if cover_image_path else _rgb(theme, "primary")

    # --- "Prepared by" — bottom-left ---
    prepared_text = f"Prepared by:  {student_name}"
    _add_textbox(
        slide,
        left=Inches(0.6), top=H - Inches(1.05),
        width=Inches(7), height=Inches(0.55),
        text=prepared_text,
        font_name=FONT_BODY, font_size=Pt(15),
        bold=True,
        color=WHITE if cover_image_path else _rgb(theme, "primary"),
        align=PP_ALIGN.LEFT
    )


# ---------------------------------------------------------------------------
# Content Slide
# ---------------------------------------------------------------------------

def _build_content_slide(prs, slide_info: dict, slide_num: int,
                         student_name: str, student_id: str, theme: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])   # blank layout

    W, H = prs.slide_width, prs.slide_height
    primary   = _rgb(theme, "primary")
    secondary = _rgb(theme, "secondary")
    accent    = _rgb(theme, "accent")
    bg        = _rgb(theme, "bg")
    text_col  = _rgb(theme, "text")
    footer_c  = _rgb(theme, "footer")

    # --- Background ---
    _add_rect(slide, 0, 0, W, H, bg)

    # --- Top header bar ---
    header_h = Inches(0.75)
    _add_rect(slide, 0, 0, W, header_h, primary)

    # --- Thin accent rule below header ---
    _add_rect(slide, 0, header_h, W, Inches(0.06), accent)

    # --- Slide title (in header bar) ---
    title_text = slide_info.get("title", f"Slide {slide_num}").upper()
    _add_textbox(
        slide,
        left=Inches(0.35), top=Inches(0.05),
        width=W - Inches(1.5), height=header_h - Inches(0.1),
        text=title_text,
        font_name=FONT_TITLE, font_size=Pt(24),
        bold=True, color=WHITE,
        align=PP_ALIGN.LEFT
    )

    # --- Page number removed as requested ---

    # --- Image (right column) ---
    has_image = False
    image_path = slide_info.get("image_path")
    if image_path and os.path.exists(image_path):
        try:
            img_left  = W - Inches(3.6)
            img_top   = header_h + Inches(0.2)
            img_w     = Inches(3.3)
            img_h     = H - header_h - Inches(0.65)
            slide.shapes.add_picture(image_path, img_left, img_top,
                                     width=img_w, height=img_h)
            has_image = True
        except Exception as e:
            print(f"[Generator] Image error: {e}")

    # --- Bullet text box ---
    body_left  = Inches(0.35)
    body_top   = header_h + Inches(0.18)
    body_w     = (W - Inches(4.1)) if has_image else (W - Inches(0.7))
    body_h     = H - header_h - Inches(0.75) # Extra safety room above footer

    txb = slide.shapes.add_textbox(body_left, body_top, body_w, body_h)
    tf  = txb.text_frame
    tf.word_wrap = True

    bullets = slide_info.get("bullets", [])
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[i] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(4)
        p.space_after  = Pt(6)

        # Arrow bullet character
        run = p.add_run()
        run.text           = f"•  {str(bullet).strip()}"
        run.font.name      = FONT_BODY
        run.font.size      = Pt(16)
        run.font.bold      = False
        run.font.color.rgb = text_col
        p.line_spacing     = 1.15
        p.space_after      = Pt(10)

    # --- Left accent strip (2px) ---
    _add_rect(slide, 0, 0, Inches(0.12), H, primary)

    # --- Footer bar ---
    footer_h   = Inches(0.42)
    footer_top = H - footer_h
    _add_rect(slide, 0, footer_top, W, footer_h, primary)

    _add_textbox(
        slide,
        left=Inches(0.3), top=H - Inches(0.38),
        width=Inches(7), height=Inches(0.3),
        text=f"Prepared by: {student_name}",
        font_name=FONT_FOOTER, font_size=Pt(10),
        bold=False, color=WHITE,
        align=PP_ALIGN.LEFT
    )

    # --- Speaker Notes ---
    notes = slide_info.get("speaker_notes", "")
    if notes:
        slide.notes_slide.notes_text_frame.text = notes


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------

def generate_pptx(slides_data: list, output_path: str,
                  student_name: str = "Student", student_id: str = "",
                  theme: str = "blue", cover_page_image=None,
                  first_slide_title: str = None) -> str:

    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    t = THEMES.get(theme, THEMES["blue"])

    # Determine the cover title
    if not first_slide_title:
        first_slide_title = (slides_data[0].get("title", "Presentation")
                             if slides_data else "Presentation")

    # Build cover slide
    _build_cover_slide(prs, first_slide_title, student_name, student_id,
                       t, cover_page_image)

    # Build content slides (skip index 0 if it is a bare cover placeholder)
    content_slides = slides_data
    if content_slides and not content_slides[0].get("bullets"):
        content_slides = content_slides[1:]  # Drop empty cover placeholder

    for num, slide_info in enumerate(content_slides, start=1):
        _build_content_slide(prs, slide_info, num, student_name, student_id, t)

    prs.save(output_path)
    print(f"[Generator] Saved PPTX → {output_path}  ({len(content_slides)} content slides)")
    return output_path


# ---------------------------------------------------------------------------
# HTML / Reveal.js Export
# ---------------------------------------------------------------------------

def _bullet_html(bullets: list) -> str:
    items = "".join(f"<li>{html_lib.escape(str(b))}</li>" for b in bullets)
    return f"<ul>{items}</ul>"


def generate_html(slides_data: list, output_path: str,
                  student_name: str = "Student", student_id: str = "") -> str:
    slides_html = ""
    for i, s in enumerate(slides_data):
        title   = html_lib.escape(s.get("title", f"Slide {i+1}"))
        bullets = _bullet_html(s.get("bullets", []))
        img_tag = ""
        img_path = s.get("image_path")
        if img_path and os.path.exists(img_path):
            with open(img_path, "rb") as f:
                b64 = base64.b64encode(f.read()).decode()
            ext = os.path.splitext(img_path)[1].strip(".") or "jpg"
            img_tag = f'<img src="data:image/{ext};base64,{b64}" class="slide-img" />'

        slides_html += f"""
        <section>
          <div class="slide-inner">
            <div class="slide-body">
              <h2>{title}</h2>
              {bullets}
            </div>
            {img_tag}
          </div>
          <div class="slide-footer">Prepared by: {html_lib.escape(student_name)}</div>
        </section>"""

    html_content = f"""<!DOCTYPE html>
<html lang="en">
<head>
  <meta charset="UTF-8">
  <meta name="viewport" content="width=device-width, initial-scale=1.0">
  <title>Presentation — {html_lib.escape(student_name)}</title>
  <link rel="stylesheet" href="https://cdnjs.cloudflare.com/ajax/libs/reveal.js/4.3.1/reset.min.css">
  <link rel="stylesheet" href="https://cdnjs.cloudflare.com/ajax/libs/reveal.js/4.3.1/reveal.min.css">
  <link rel="stylesheet" href="https://cdnjs.cloudflare.com/ajax/libs/reveal.js/4.3.1/theme/night.min.css">
  <style>
    :root {{
      --blue: #1f497d;
      --gold: #ffc000;
      --light: #f5f8ff;
    }}
    .reveal h2 {{
      font-family: 'Calibri', sans-serif;
      font-size: 1.6em;
      color: var(--gold);
      text-transform: uppercase;
      letter-spacing: 0.03em;
      margin-bottom: 0.5em;
    }}
    .reveal ul {{
      font-family: 'Calibri', sans-serif;
      font-size: 0.85em;
      line-height: 1.7;
      text-align: left;
      margin-left: 1.2em;
    }}
    .reveal li {{ margin-bottom: 0.35em; }}
    .slide-inner {{
      display: flex;
      gap: 1.5em;
      align-items: flex-start;
    }}
    .slide-body {{ flex: 1; }}
    .slide-img {{
      width: 280px;
      height: 200px;
      object-fit: cover;
      border-radius: 8px;
      box-shadow: 0 4px 18px rgba(0,0,0,0.5);
    }}
    .slide-footer {{
      position: absolute;
      bottom: 10px;
      left: 30px;
      font-size: 0.55em;
      opacity: 0.7;
      font-family: 'Calibri', sans-serif;
    }}
  </style>
</head>
<body>
<div class="reveal">
  <div class="slides">
    {slides_html}
  </div>
</div>
<script src="https://cdnjs.cloudflare.com/ajax/libs/reveal.js/4.3.1/reveal.min.js"></script>
<script>
  Reveal.initialize({{
    hash: true, transition: 'convex',
    backgroundTransition: 'fade', slideNumber: true
  }});
</script>
</body>
</html>"""

    with open(output_path, "w", encoding="utf-8") as f:
        f.write(html_content)
    print(f"[Generator] Saved HTML → {output_path}")
    return output_path


# ---------------------------------------------------------------------------
# Utilities
# ---------------------------------------------------------------------------

def get_image_base64(image_path: str):
    if not image_path or not os.path.exists(image_path):
        return None
    with open(image_path, "rb") as f:
        b64 = base64.b64encode(f.read()).decode()
    ext = os.path.splitext(image_path)[1].strip(".") or "jpg"
    return f"data:image/{ext};base64,{b64}"
